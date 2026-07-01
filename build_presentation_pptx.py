"""
Generate presentation.pptx -- a PowerPoint version of presentation.html, same
12-slide content (GeoChron paper adaptation walkthrough). Run once; re-run
after editing the SLIDES list below to regenerate.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

DARK_BG   = RGBColor(0x1B, 0x1E, 0x24)
CARD_BG   = RGBColor(0x26, 0x2A, 0x33)
TEXT      = RGBColor(0xEA, 0xEA, 0xEA)
BODY      = RGBColor(0xCF, 0xD3, 0xDA)
KICKER    = RGBColor(0x6B, 0x72, 0x80)
BLUE      = RGBColor(0x7F, 0xD1, 0xFF)   # paper-original
GREEN     = RGBColor(0x7F, 0xFF, 0xA0)   # adaptation
ORANGE    = RGBColor(0xFF, 0x9D, 0x7F)   # difference / not-in-paper

TAG_COLORS = {"paper": BLUE, "adapt": GREEN, "diff": ORANGE}

# Each slide: kicker, title, list of (tag, text) card paragraphs, optional table (header, rows)
SLIDES = [
    dict(kicker="GeoChron × Texas PMIS",
         title="Adapting a Storyline visualization paper to real pavement data",
         body=[(None, "What the paper proposed, what we kept, what we had to change, and why — "
                      "a walkthrough of the cohort-based pavement trajectory pipeline.")]),

    dict(kicker="1 · The paper's core idea", title="GeoChron's Storyline framework",
         body=[("paper", "Treat each spatiotemporal (ST) series (an entity over time) as a "
                          "\"character,\" and a group of series moving together as a Storyline "
                          "session — the same analogy classic Storyline visualizations use for "
                          "characters interacting in a narrative."),
               ("paper", "Pipeline per time slice: (1) split time into slices, (2) correlate every "
                         "pair of series within a slice, (3) intersect correlation with spatial "
                         "proximity to build a relation network, (4) run Louvain community "
                         "detection on that network — each community = one session, (5) lay "
                         "sessions out with a StoryFlow-style layout, (6) drill down with an "
                         "EvoLens showing per-entity trends and a normalized \"trend motif.\"")]),

    dict(kicker="2 · Our adaptation target", title="Texas PMIS pavement condition data",
         body=[(None, "The paper's data: dense, regularly-sampled ST series; clean single-day time "
                      "slices; coordinates available for every entity."),
               (None, "Our data (PMIS): ~14,810 half-mile road segments, 1996–2024; annual condition "
                      "scores, sparse and gappy; begin_dfo/end_dfo coordinates empty for our data."),
               (None, "Same underlying idea, a very different data reality — this drove every "
                      "adaptation decision that follows.")]),

    dict(kicker="3 · Adaptation: time slice → sliding window",
         title="Step 4: multi-year sliding windows replace single-day slices",
         body=[("diff", "Paper: a time slice is a single day; a time window wraps three "
                        "consecutive slices for cross-slice smoothing."),
               ("adapt", "Ours: annual sampling can't support a single-day slice. We use one "
                         "W=5-year sliding window, stepped by 1 year (1996–2024 → 25 windows). "
                         "The window overlap gives us the paper's cross-slice smoothing \"for "
                         "free.\""),
               (None, "Called out in procedure.md as the one non-negotiable adaptation.")]),

    dict(kicker="4 · Adaptation: gaps", title="Step 3: preserve gaps, never interpolate",
         body=[("diff", "With sparse annual sampling, an interpolated point can be a large "
                        "fraction of a segment's signal — enough to manufacture a trend the "
                        "pavement never had."),
               ("adapt", "Our rule: leave missing years empty in the section×year matrix. Handle "
                         "gaps only at correlation time (pairwise-complete Pearson, minimum "
                         "overlap required) — never by fabricating values.")]),

    dict(kicker="5 · Adaptation: proximity without coordinates",
         title="Step 8–9: (roadbed, county) instead of geographic radius",
         table=(["Paper", "Ours"],
                [["Same highway within a reference-marker distance, OR same district / within a "
                  "haversine radius (cross-route corridor effects)",
                  "Same roadbed AND county — a proxy, since real lat/long isn't populated in this "
                  "dataset"]]),
         body=[(None, "The biggest known simplification in the pipeline: cohorts can only be found "
                      "within one highway/county pair, never spanning two different but nearby "
                      "highways. Flagged in PROGRESS.md as the first thing to revisit if real "
                      "coordinates become available.")]),

    dict(kicker="6 · Straight reuse (no change needed)",
         title="What we kept exactly as the paper describes",
         body=[("paper", "Step 6–7 — correlation + thresholding: pairwise-complete Pearson per "
                         "window, reduced to a binary yes/no edge at THR=0.7."),
               ("paper", "Step 10 — Louvain community detection, run on the intersection of "
                         "correlation and proximity edges; identical mechanism to the paper."),
               ("paper", "Step 11–12 — session filtering + cross-window tracking by largest "
                         "membership overlap, producing the Storyline's continuity curves.")]),

    dict(kicker="6b · How much survives the pipeline",
         title="The funnel: from 41,314 raw sections to storyline_data.json",
         table=(["Stage", "Kept", "Filtered out"],
                [["Step 2 — needs ≥10 years of real data", "14,812 sections", "26,502 dropped (64%)"],
                 ["Step 6+7 — correlated pairs (r>0.7), all windows", "82,574,973 edges", "— (already a small slice of all possible pairs)"],
                 ["Step 8+9 — intersect with (roadbed, county) proximity", "506,905 edges", "82,068,068 dropped (99.4%)"],
                 ["Step 10 — raw Louvain communities, all windows", "118,675 communities", "—"],
                 ["Step 11 — sessions ≥ size 5 (or rescued by a neighbor window)", "24,408 sessions", "94,267 dropped (79%)"]]),
         body=[(None, "Biggest single filter: Step 8/9's spatial proximity rule discards 99.4% of "
                      "correlated pairs for not sharing a highway+county. Second biggest: Step 2's "
                      "≥10-year data requirement removes 64% of sections statewide before the "
                      "pipeline even starts."),
               (None, "Net: 14,810 of 41,314 raw sections (36%) ever appear in storyline_data.json "
                      "at all — and even among those, many appear only as unaffiliated singletons "
                      "in a given window rather than in a kept session.")]),

    dict(kicker="7 · Adaptation: normalization, deferred",
         title="Step 5: skipped for correlation, used later for EvoLens",
         body=[("adapt", "Pearson correlation is provably invariant to shift/scale, so Step 5's "
                         "normalization is mathematically redundant for the correlation math — we "
                         "skip it there and use raw scores."),
               ("adapt", "Normalization (z-score) is applied live, on demand, in the EvoLens trend "
                         "motif — where the paper genuinely needs it.")]),

    dict(kicker="8 · Visualization — Level 1",
         title="Storyline rendering: window = span, not a point",
         body=[("diff", "Early prototype (wrong): each window drawn as a single point; "
                        "single-window memberships rendered as meaningless dots; anchor/"
                        "interpolation logic produced artificial spikes into empty space."),
               ("adapt", "Current, paper-faithful: each window is a full-width bar (a span), "
                         "windows connected by a thinner Bézier curve. Layout uses per-window "
                         "barycenter-sweep ordering with no anchor points at all.")]),

    dict(kicker="9 · Visualization — engineering, not paper-driven",
         title="WebGL rendering (a practical addition, not from the paper)",
         body=[("diff", "The paper doesn't specify a rendering technology. Our first D3/SVG "
                        "implementation became too slow at ~14,810 segments × 25 windows."),
               ("adapt", "Our addition: hybrid renderer — WebGL for bulk bar/connector geometry "
                         "(chunked draw calls to dodge a GPU driver limit), plus a transparent 2D "
                         "canvas overlay for text, hover highlighting, and labels.")]),

    dict(kicker="10 · Visualization — Level 2",
         title="EvoLens drill-down: brush-select, not click",
         body=[("paper", "The paper's actual mechanism: brush a rectangle on the Storyline to "
                         "select a time period + sessions, opening a lens with per-entity line "
                         "charts of the real trend, plus a trend motif (z-score within the brushed "
                         "range → lower quartile / median / upper quartile band)."),
               (None, "We initially built this as a click interaction; corrected to brush-select "
                      "once verified against the paper text.")]),

    dict(kicker="11 · Display-only fix", title="Grouping re-keyed to (roadbed, county)",
         body=[(None, "Sessions were already correctly scoped to one (roadbed, county) pair "
                      "(Step 8/9). But the Storyline's visual bands were originally grouped by "
                      "roadbed alone, mixing counties a highway passes through."),
               ("adapt", "Fix: re-keyed export_storyline_data.py's band grouping to "
                         "(roadbed, county) tuples: 722 → 1,112 bands. No session/track validity "
                         "changed — only the display grouping.")]),

    dict(kicker="12 · Summary", title="Same paper, different data reality",
         table=(["Kept exactly as the paper", "Adapted for our data", "Not in the paper (engineering)"],
                [["Correlation → threshold → Louvain → session filter → cross-window tracking → "
                  "brush-select EvoLens → trend motif",
                  "Time slice → 5-yr sliding window · gaps preserved not interpolated · proximity "
                  "via (roadbed, county) proxy · normalization deferred to EvoLens",
                  "WebGL/canvas hybrid rendering · JSON data-export + static JS front-end architecture"]]),
         body=[(None, "Open items: Step 13 (treatment-event detection), Step 14 (validation "
                      "against ground truth), Step 15 (formal run-scope decision), the paper's "
                      "linked geographic map (blocked on missing coordinates), and independent "
                      "confirmation of the score-direction convention.")]),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]


def set_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = DARK_BG


def add_text(slide, left, top, width, height, text, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.name = "Segoe UI"
    return box


def add_card(slide, left, top, width, height, tag, text):
    box = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE.RECTANGLE = 1
    box.fill.solid()
    box.fill.fore_color.rgb = CARD_BG
    box.line.color.rgb = TAG_COLORS.get(tag, RGBColor(0x38, 0x3D, 0x47))
    box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.18)
    tf.margin_right = Inches(0.18)
    tf.margin_top = Inches(0.12)
    tf.margin_bottom = Inches(0.12)
    if tag:
        p0 = tf.paragraphs[0]
        r0 = p0.add_run()
        r0.text = {"paper": "PAPER", "adapt": "ADAPTED", "diff": "DIFFERENT"}[tag]
        r0.font.size = Pt(11)
        r0.font.bold = True
        r0.font.color.rgb = TAG_COLORS[tag]
        p1 = tf.add_paragraph()
    else:
        p1 = tf.paragraphs[0]
    r1 = p1.add_run()
    r1.text = text
    r1.font.size = Pt(15)
    r1.font.color.rgb = BODY


def add_table(slide, left, top, width, height, header, rows):
    n_rows, n_cols = len(rows) + 1, len(header)
    gtable = slide.shapes.add_table(n_rows, n_cols, left, top, width, height).table
    for c, h in enumerate(header):
        cell = gtable.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = CARD_BG
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.color.rgb = BLUE
    for r_i, row in enumerate(rows, start=1):
        for c_i, val in enumerate(row):
            cell = gtable.cell(r_i, c_i)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_BG
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(13)
                    run.font.color.rgb = BODY


for spec in SLIDES:
    slide = prs.slides.add_slide(blank)
    set_bg(slide)
    add_text(slide, Inches(0.7), Inches(0.4), Inches(12), Inches(0.4),
             spec["kicker"], 13, KICKER, bold=True)
    add_text(slide, Inches(0.7), Inches(0.75), Inches(12), Inches(1.1),
             spec["title"], 30, RGBColor(0xFF, 0xFF, 0xFF), bold=True)

    y = Inches(2.0)
    if "table" in spec:
        header, rows = spec["table"]
        table_h = Inches(1.1 + 0.5 * len(rows))
        add_table(slide, Inches(0.7), y, Inches(12), table_h, header, rows)
        y = y + table_h + Inches(0.25)

    for tag, text in spec["body"]:
        card_h = Inches(0.35 + 0.023 * len(text))
        add_card(slide, Inches(0.7), y, Inches(12), card_h, tag, text)
        y = y + card_h + Inches(0.18)

prs.save("presentation.pptx")
print(f"wrote presentation.pptx ({len(SLIDES)} slides)")
